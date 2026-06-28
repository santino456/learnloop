from __future__ import annotations

import hashlib
import json
import re
import shutil
import warnings
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from .model import LearnLoopError


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".md", ".markdown", ".txt"}
CROP_RENDER_SCALE = 4
CAPTION_RE = re.compile(
    r"^\s*((?:fig(?:ure)?\.?|图|table|表)\s*[\dIVXLC]+[a-zA-Z]?)\s*[:.\-：]\s*(.*)",
    re.I,
)


@dataclass
class IngestResult:
    source: Path
    material_dir: Path
    material_json: Path
    chunks_jsonl: Path
    figures_json: Path | None = None
    figures_md: Path | None = None
    figures: int = 0
    chunks: int = 0


def ingest_material(
    source: Path,
    course_dir: Path,
    *,
    backend: str = "auto",
    overwrite: bool = False,
) -> IngestResult:
    """Convert a raw source file into a LearnLoop material pack."""
    source = source.expanduser().resolve()
    course_dir = course_dir.expanduser().resolve()
    if not source.exists():
        raise LearnLoopError(f"Source file does not exist: {source}")
    if not course_dir.exists():
        raise LearnLoopError(f"Course directory does not exist: {course_dir}")

    ext = source.suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise LearnLoopError(
            f"Unsupported source type: {ext or '(none)'}. "
            "Supported: PDF, DOCX, PPTX, Markdown, text."
        )

    material_dir = course_dir / ".learnloop" / "materials" / _safe_stem(source)
    if material_dir.exists() and overwrite:
        shutil.rmtree(material_dir)
    material_dir.mkdir(parents=True, exist_ok=True)
    (course_dir / "assets").mkdir(exist_ok=True)

    if ext == ".pdf":
        chunks, figures, metadata = _ingest_pdf(source, course_dir, material_dir, backend)
        selected_backend = metadata["backend"]
    elif ext == ".docx":
        chunks, figures, metadata = _ingest_docx(source, material_dir, backend)
        selected_backend = metadata["backend"]
    elif ext == ".pptx":
        chunks, figures, metadata = _ingest_pptx(source, material_dir, backend)
        selected_backend = metadata["backend"]
    else:
        chunks, figures, metadata = _ingest_text(source, material_dir, backend)
        selected_backend = metadata["backend"]

    chunks_jsonl = material_dir / "chunks.jsonl"
    _write_jsonl(chunks_jsonl, chunks)

    figures_json: Path | None = None
    figures_md: Path | None = None
    if figures:
        figures_json = material_dir / "figures.json"
        figures_json.write_text(
            json.dumps({"figures": figures}, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
        )
        figures_md = material_dir / "figures.md"
        figures_md.write_text(_figures_markdown(figures), encoding="utf-8")

    material = {
        "schema_version": 1,
        "source": {
            "path": str(source),
            "filename": source.name,
            "type": ext.lstrip("."),
        },
        "backend": selected_backend,
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "outputs": {
            "chunks_jsonl": "chunks.jsonl",
            "figures_json": "figures.json" if figures_json else None,
            "figures_md": "figures.md" if figures_md else None,
        },
        "stats": {
            "chunks": len(chunks),
            "figures": len(figures),
        },
        "metadata": metadata,
    }
    material_json = material_dir / "material.json"
    material_json.write_text(
        json.dumps(material, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    _register_source(course_dir, source, material_dir, material["stats"])

    return IngestResult(
        source=source,
        material_dir=material_dir,
        material_json=material_json,
        chunks_jsonl=chunks_jsonl,
        figures_json=figures_json,
        figures_md=figures_md,
        figures=len(figures),
        chunks=len(chunks),
    )


def _ingest_pdf(
    source: Path, course_dir: Path, material_dir: Path, backend: str
) -> tuple[list[dict[str, Any]], list[dict[str, Any]], dict[str, Any]]:
    if backend not in {"auto", "pymupdf"}:
        raise LearnLoopError(
            f"PDF backend '{backend}' is not available in the base install. "
            "Use 'pymupdf' or install an advanced extractor outside LearnLoop."
        )
    try:
        import fitz  # type: ignore
    except ImportError as exc:  # pragma: no cover - depends on environment
        raise LearnLoopError(
            "PDF ingest requires PyMuPDF. Install it with: python3 -m pip install pymupdf"
        ) from exc

    doc = fitz.open(source)
    chunks: list[dict[str, Any]] = []
    figure_candidates: list[dict[str, Any]] = []
    source_id = _safe_stem(source)
    assets_dir = course_dir / "assets"

    for page_index, page in enumerate(doc, start=1):
        text = page.get_text("text").strip()
        for chunk_index, paragraph in enumerate(_split_paragraphs(text), start=1):
            chunks.append(
                {
                    "id": f"{source_id}:p{page_index}:c{chunk_index}",
                    "source": source.name,
                    "page": page_index,
                    "kind": "text",
                    "text": paragraph,
                }
            )

        caption_blocks = _caption_blocks(page)
        text_blocks = _text_blocks(page)
        tables = _table_blocks(page)
        table_rects = [item["bbox"] for item in tables]
        visual_rects = _visual_rects(page)

        for caption in caption_blocks:
            if caption["kind"] == "table":
                rect = _infer_table_rect(
                    page, caption["bbox"], tables, text_blocks=text_blocks
                )
            else:
                rect = _infer_figure_rect(
                    page.rect,
                    caption["bbox"],
                    visual_rects,
                    text_blocks=text_blocks,
                    table_rects=table_rects,
                    kind=caption["kind"],
                )
            if rect is None:
                continue
            figure_candidates.append(
                {
                    "page": page_index,
                    "caption": caption,
                    "rect": rect,
                }
            )

    figures: list[dict[str, Any]] = []
    seen_hashes: dict[str, int] = {}
    for candidate in figure_candidates:
        page_index = candidate["page"]
        page = doc[page_index - 1]
        rect = candidate["rect"]
        caption = candidate["caption"]
        pix = page.get_pixmap(matrix=fitz.Matrix(CROP_RENDER_SCALE, CROP_RENDER_SCALE), clip=rect, alpha=False)
        image_hash = hashlib.sha256(pix.tobytes()).hexdigest()
        if image_hash in seen_hashes:
            continue

        global_index = len(figures) + 1
        asset_name = _label_asset_name(source_id, caption["label"], caption["kind"], global_index)
        asset_path = assets_dir / asset_name
        pix.save(asset_path)

        figure_record: dict[str, Any] = {
            "id": f"{source_id}:{asset_name[:-4]}",
            "source": source.name,
            "page": page_index,
            "label": caption["label"],
            "caption": caption["text"],
            "kind": caption["kind"],
            "asset": f"assets/{asset_name}",
            "bbox": [round(value, 2) for value in rect],
            "method": "caption-nearby-crop",
        }
        if caption["kind"] == "table":
            table_md = _extract_table_markdown(page, rect)
            if table_md:
                figure_record["table_markdown"] = table_md

        figures.append(figure_record)
        seen_hashes[image_hash] = global_index

    metadata = {
        "backend": "pymupdf",
        "pages": doc.page_count,
        "title": doc.metadata.get("title") or "",
        "author": doc.metadata.get("author") or "",
        "figure_extraction": "caption-nearby-crop",
    }
    doc.close()
    return chunks, figures, metadata


def _ingest_docx(
    source: Path, _material_dir: Path, backend: str
) -> tuple[list[dict[str, Any]], list[dict[str, Any]], dict[str, Any]]:
    if backend not in {"auto", "python-docx"}:
        raise LearnLoopError("DOCX ingest currently supports backend 'python-docx'.")
    try:
        import docx  # type: ignore
    except ImportError as exc:  # pragma: no cover - depends on environment
        raise LearnLoopError(
            "DOCX ingest requires python-docx. Install it with: python3 -m pip install python-docx"
        ) from exc

    document = docx.Document(str(source))
    source_id = _safe_stem(source)
    chunks = []
    for index, paragraph in enumerate(document.paragraphs, start=1):
        text = paragraph.text.strip()
        if not text:
            continue
        style = paragraph.style.name if paragraph.style else ""
        chunks.append(
            {
                "id": f"{source_id}:p{index}",
                "source": source.name,
                "kind": "heading" if style.lower().startswith("heading") else "text",
                "style": style,
                "text": text,
            }
        )
    return chunks, [], {"backend": "python-docx", "paragraphs": len(chunks)}


def _ingest_pptx(
    source: Path, _material_dir: Path, backend: str
) -> tuple[list[dict[str, Any]], list[dict[str, Any]], dict[str, Any]]:
    if backend not in {"auto", "python-pptx"}:
        raise LearnLoopError("PPTX ingest currently supports backend 'python-pptx'.")
    try:
        import pptx  # type: ignore
    except ImportError as exc:  # pragma: no cover - depends on environment
        raise LearnLoopError(
            "PPTX ingest requires python-pptx. Install it with: python3 -m pip install python-pptx"
        ) from exc

    presentation = pptx.Presentation(str(source))
    source_id = _safe_stem(source)
    chunks = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    texts.append(text)
        if texts:
            chunks.append(
                {
                    "id": f"{source_id}:s{slide_index}",
                    "source": source.name,
                    "slide": slide_index,
                    "kind": "slide",
                    "text": "\n".join(texts),
                }
            )
    return chunks, [], {"backend": "python-pptx", "slides": len(presentation.slides)}


def _ingest_text(
    source: Path, _material_dir: Path, backend: str
) -> tuple[list[dict[str, Any]], list[dict[str, Any]], dict[str, Any]]:
    if backend not in {"auto", "plain-text"}:
        raise LearnLoopError("Text ingest currently supports backend 'plain-text'.")
    source_id = _safe_stem(source)
    text = source.read_text(encoding="utf-8")
    chunks = [
        {
            "id": f"{source_id}:c{index}",
            "source": source.name,
            "kind": "text",
            "text": paragraph,
        }
        for index, paragraph in enumerate(_split_paragraphs(text), start=1)
    ]
    return chunks, [], {"backend": "plain-text"}


def _caption_blocks(page: Any) -> list[dict[str, Any]]:
    """Extract and classify caption blocks, merging multi-line captions."""
    raw_blocks: list[dict[str, Any]] = []
    for block in page.get_text("blocks"):
        if len(block) < 5:
            continue
        x0, y0, x1, y1, text, *_ = block
        text = _clean_text(text)
        if not text:
            continue
        raw_blocks.append(
            {
                "bbox": page.rect.__class__(x0, y0, x1, y1),
                "text": text,
            }
        )
    raw_blocks.sort(key=lambda b: (b["bbox"].y0, b["bbox"].x0))

    captions: list[dict[str, Any]] = []
    i = 0
    while i < len(raw_blocks):
        block = raw_blocks[i]
        match = CAPTION_RE.match(block["text"])
        if not match:
            i += 1
            continue

        label = match.group(1).strip()
        caption_text = block["text"]
        bbox = block["bbox"]
        kind = _caption_kind(label)
        j = i + 1
        while j < len(raw_blocks):
            nxt = raw_blocks[j]
            prev = raw_blocks[j - 1]
            gap = nxt["bbox"].y0 - prev["bbox"].y1
            same_line = abs(nxt["bbox"].y0 - bbox.y0) < 4 or gap < 4
            if not same_line:
                break
            overlap = min(nxt["bbox"].x1, bbox.x1) - max(nxt["bbox"].x0, bbox.x0)
            if overlap < -page.rect.width * 0.3:
                break
            if CAPTION_RE.match(nxt["text"]):
                break
            caption_text += " " + nxt["text"]
            bbox |= nxt["bbox"]
            j += 1

        captions.append(
            {
                "bbox": bbox,
                "label": label,
                "text": caption_text,
                "kind": kind,
            }
        )
        i = j if j > i else i + 1

    return captions


def _caption_kind(label: str) -> str:
    """Classify a caption label as figure, table, or subfigure."""
    low = label.lower()
    if "table" in low or low.startswith("表"):
        return "table"
    # Labels like Fig 1a, Fig 2(b), 图 3-A
    if re.search(r"[\dIVXLC]+\s*[\(\[]?[a-zA-Z][\)\]]?$", label, re.I):
        return "subfigure"
    return "figure"


def _visual_rects(page: Any, exclude_rects: list[Any] | None = None) -> list[Any]:
    """Collect image and drawing rectangles on a page, warning on errors.

    The ``exclude_rects`` parameter is kept for backwards compatibility but no
    longer filters rectangles here; callers filter visual candidates against
    tables or text blocks as needed for their caption type.
    """
    rects: list[Any] = []
    page_area = page.rect.width * page.rect.height
    max_rect_area = page_area * 0.8
    min_rect_area = 100
    seen: set[tuple[int, int, int, int]] = set()

    for image in page.get_images(full=True):
        xref = image[0]
        try:
            for rect in page.get_image_rects(xref):
                key = (
                    round(rect.x0),
                    round(rect.y0),
                    round(rect.x1),
                    round(rect.y1),
                )
                if key in seen:
                    continue
                seen.add(key)
                rects.append(rect)
        except Exception as exc:
            warnings.warn(
                f"Could not get image rects for xref {xref} on page {page.number + 1}: {exc}"
            )
    try:
        for drawing in page.get_drawings():
            rect = drawing.get("rect")
            if rect is None:
                continue
            area = rect.width * rect.height if hasattr(rect, "width") else 0
            if not (min_rect_area <= area <= max_rect_area):
                continue
            key = (
                round(rect.x0),
                round(rect.y0),
                round(rect.x1),
                round(rect.y1),
            )
            if key in seen:
                continue
            seen.add(key)
            rects.append(rect)
    except Exception as exc:
        warnings.warn(f"Could not get drawings on page {page.number + 1}: {exc}")

    return sorted(rects, key=lambda r: (r.y0, r.x0))


def _text_blocks(page: Any) -> list[dict[str, Any]]:
    """Return text blocks with bboxes, excluding empty/whitespace blocks."""
    blocks: list[dict[str, Any]] = []
    for block in page.get_text("blocks"):
        if len(block) < 5:
            continue
        x0, y0, x1, y1, text, *_ = block
        text = _clean_text(text)
        if not text:
            continue
        blocks.append(
            {
                "bbox": page.rect.__class__(x0, y0, x1, y1),
                "text": text,
            }
        )
    return blocks


def _table_blocks(page: Any) -> list[dict[str, Any]]:
    """Return table objects with normalized bboxes.

    Filters out tiny or degenerate detections (e.g. a chart misread as a table
    with very few cells) so callers can treat the remainder as genuine tables.
    """
    tables: list[dict[str, Any]] = []
    try:
        for table in page.find_tables():
            try:
                if table.row_count < 2 or table.col_count < 2 or len(table.cells) < 6:
                    continue
                tables.append(
                    {
                        "bbox": page.rect.__class__(*table.bbox),
                        "table": table,
                    }
                )
            except Exception:
                continue
    except Exception as exc:
        warnings.warn(f"Could not find tables on page {page.number + 1}: {exc}")
    return tables


def _rect_overlap_area(a: Any, b: Any) -> float:
    x_overlap = max(0, min(a.x1, b.x1) - max(a.x0, b.x0))
    y_overlap = max(0, min(a.y1, b.y1) - max(a.y0, b.y0))
    return x_overlap * y_overlap


def _rect_overlap_ratio(a: Any, b: Any) -> float:
    area = a.width * a.height if hasattr(a, "width") else 0
    if area <= 0:
        return 0.0
    return _rect_overlap_area(a, b) / area


def _shrink_table_rect(
    rect: Any,
    text_blocks: list[dict[str, Any]],
    gap_threshold: float,
) -> Any:
    """Shrink a table bbox by cutting at the first large vertical text gap.

    This removes body text or figures that a text-strategy table finder merged
    into the table region. Overlapping blocks are merged before gap detection.
    """
    intervals: list[list[float]] = []
    for block in text_blocks:
        bbox = block["bbox"]
        if _rect_overlap_area(bbox, rect) <= 0:
            continue
        intervals.append([bbox.y0, bbox.y1])
    if not intervals:
        return rect

    intervals.sort(key=lambda pair: pair[0])
    merged: list[list[float]] = [intervals[0]]
    for y0, y1 in intervals[1:]:
        prev = merged[-1]
        if y0 <= prev[1]:
            prev[1] = max(prev[1], y1)
        else:
            merged.append([y0, y1])

    # Shrink top if there is a gap before the first text.
    if merged[0][0] > rect.y0 + gap_threshold:
        rect.y0 = merged[0][0]

    # Cut bottom at the first large vertical gap.
    for i in range(len(merged) - 1):
        gap = merged[i + 1][0] - merged[i][1]
        if gap > gap_threshold:
            rect.y1 = merged[i][1]
            break

    # Also shrink left/right to the actual text extent.
    overlapping = [b["bbox"] for b in text_blocks if _rect_overlap_area(b["bbox"], rect) > 0]
    if overlapping:
        rect.x0 = max(rect.x0, min(b.x0 for b in overlapping))
        rect.x1 = min(rect.x1, max(b.x1 for b in overlapping))
    return rect


def _infer_table_rect(
    page: Any,
    caption_rect: Any,
    tables: list[dict[str, Any]],
    text_blocks: list[dict[str, Any]] | None = None,
    margin: int = 6,
) -> Any | None:
    """Infer the bounding box of a table from its caption.

    First tries PyMuPDF's default line-based table detector. If that fails to
    find a plausible table near the caption, it falls back to the text-based
    strategy clipped to the region directly adjacent to the caption, then
    shrinks the result at the first large vertical text gap so body text or
    nearby figures are not included.
    """
    text_blocks = text_blocks or []
    page_rect = page.rect
    max_gap = page_rect.height * 0.55
    min_clearance = 4

    def nearby(rect: Any) -> bool:
        above = (
            rect.y1 <= caption_rect.y0 - min_clearance
            and rect.y1 >= caption_rect.y0 - max_gap
        )
        below = (
            rect.y0 >= caption_rect.y1 + min_clearance
            and rect.y0 <= caption_rect.y1 + max_gap
        )
        if not (above or below):
            return False
        horizontal_overlap = min(rect.x1, caption_rect.x1) - max(
            rect.x0, caption_rect.x0
        )
        return horizontal_overlap > -page_rect.width * 0.35

    line_candidates = [
        (abs(_vertical_distance(caption_rect, item["bbox"])), item["bbox"])
        for item in tables
        if nearby(item["bbox"])
    ]

    line_rect: Any | None = None
    if line_candidates:
        line_candidates.sort(key=lambda x: x[0])
        line_rect = page_rect.__class__(*line_candidates[0][1])

    # Always try the text-based strategy as well; it is more robust for tables
    # without ruling lines and lets us choose the candidate closest to caption.
    text_rect = _infer_table_rect_from_text(
        page, page_rect, caption_rect, text_blocks
    )

    table_rect: Any | None = None
    if line_rect is not None and text_rect is not None:
        table_rect = (
            line_rect
            if abs(_vertical_distance(caption_rect, line_rect))
            <= abs(_vertical_distance(caption_rect, text_rect))
            else text_rect
        )
    elif line_rect is not None:
        table_rect = line_rect
    else:
        table_rect = text_rect

    if table_rect is None:
        return None

    # Remove body text / figures merged into a text-strategy table bbox.
    gap_threshold = page_rect.height * 0.02
    table_rect = _shrink_table_rect(table_rect, text_blocks, gap_threshold)

    figure_rect = page_rect.__class__(
        max(page_rect.x0, table_rect.x0 - margin),
        max(page_rect.y0, table_rect.y0 - margin),
        min(page_rect.x1, table_rect.x1 + margin),
        min(page_rect.y1, table_rect.y1 + margin),
    )

    if table_rect.y1 <= caption_rect.y0:
        figure_rect.y1 = min(figure_rect.y1, caption_rect.y0 - min_clearance)
    else:
        figure_rect.y0 = max(figure_rect.y0, caption_rect.y1 + min_clearance)

    return _clip_rect(figure_rect, page_rect)


def _infer_table_rect_from_text(
    page: Any,
    page_rect: Any,
    caption_rect: Any,
    text_blocks: list[dict[str, Any]],
) -> Any | None:
    """Use PyMuPDF's text-based table finder in the caption-adjacent region."""
    max_height = page_rect.height * 0.5
    min_height = 48
    min_clearance = 4

    above = page_rect.__class__(
        page_rect.x0,
        max(page_rect.y0, caption_rect.y0 - max_height),
        page_rect.x1,
        caption_rect.y0 - min_clearance,
    )
    below = page_rect.__class__(
        page_rect.x0,
        caption_rect.y1 + min_clearance,
        page_rect.x1,
        min(page_rect.y1, caption_rect.y1 + max_height),
    )

    best: Any | None = None
    best_score = 0.0
    for region in (above, below):
        if region.y1 - region.y0 < min_height:
            continue
        try:
            finder = page.find_tables(
                clip=region,
                vertical_strategy="text",
                horizontal_strategy="text",
            )
            for table in finder:
                tbbox = page_rect.__class__(*table.bbox)
                overlap = _rect_overlap_area(tbbox, region)
                cells = len(table.cells)
                if cells < 6:
                    continue
                # Prefer tables that are close to the caption and have many cells.
                score = overlap + cells * 10
                if score > best_score:
                    best_score = score
                    best = tbbox
        except Exception:
            continue

    return best


def _vertical_distance(a: Any, b: Any) -> float:
    if b.y1 <= a.y0:
        return a.y0 - b.y1
    if b.y0 >= a.y1:
        return b.y0 - a.y1
    return 0.0


def _text_overlap_ratio(rect: Any, text_blocks: list[dict[str, Any]]) -> float:
    if not text_blocks:
        return 0.0
    overlap = sum(_rect_overlap_area(rect, block["bbox"]) for block in text_blocks)
    area = rect.width * rect.height
    return overlap / area if area > 0 else 1.0


def _infer_figure_rect(
    page_rect: Any,
    caption_rect: Any,
    visual_rects: list[Any],
    text_blocks: list[dict[str, Any]] | None = None,
    table_rects: list[Any] | None = None,
    kind: str = "figure",
    margin: int = 24,
) -> Any | None:
    """Infer the bounding box of a figure near a caption.

    Uses spatial clustering of nearby visual elements and adapts to captions
    placed above or below the content.  The ``table_rects`` parameter is kept
    for backwards compatibility but no longer filters visual candidates here.
    """
    text_blocks = text_blocks or []
    table_rects = table_rects or []

    if not visual_rects:
        return _fallback_figure_rect(page_rect, caption_rect, margin, text_blocks)

    max_gap = page_rect.height * 0.55
    min_clearance = 4
    candidates = []

    for rect in visual_rects:
        above = (
            rect.y1 <= caption_rect.y0 - min_clearance
            and rect.y1 >= caption_rect.y0 - max_gap
        )
        below = (
            rect.y0 >= caption_rect.y1 + min_clearance
            and rect.y0 <= caption_rect.y1 + max_gap
        )
        if not (above or below):
            continue
        horizontal_overlap = min(rect.x1, caption_rect.x1) - max(
            rect.x0, caption_rect.x0
        )
        if horizontal_overlap <= -page_rect.width * 0.35:
            continue
        candidates.append(rect)

    if not candidates:
        return _fallback_figure_rect(page_rect, caption_rect, margin, text_blocks)

    clusters = _cluster_rects(candidates)

    def caption_distance(cluster: Any) -> float:
        if cluster.y1 <= caption_rect.y0:
            return caption_rect.y0 - cluster.y1
        return cluster.y0 - caption_rect.y1

    clusters.sort(key=lambda c: (caption_distance(c), -(c.width * c.height)))
    cluster = clusters[0]

    if kind == "subfigure" and len(clusters) > 1:
        sub_clusters = _split_subfigure_clusters(clusters)
        if len(sub_clusters) > 1:
            cx = (caption_rect.x0 + caption_rect.x1) / 2
            cluster = min(
                sub_clusters, key=lambda c: abs((c.x0 + c.x1) / 2 - cx)
            )

    figure_rect = page_rect.__class__(
        max(page_rect.x0, cluster.x0 - margin),
        max(page_rect.y0, cluster.y0 - 6),
        min(page_rect.x1, cluster.x1 + margin),
        min(page_rect.y1, cluster.y1 + 6),
    )

    if cluster.y1 <= caption_rect.y0:
        figure_rect.y1 = min(figure_rect.y1, caption_rect.y0 - min_clearance)
    else:
        figure_rect.y0 = max(figure_rect.y0, caption_rect.y1 + min_clearance)

    return _clip_rect(figure_rect, page_rect)


def _fallback_figure_rect(
    page_rect: Any,
    caption_rect: Any,
    margin: int,
    text_blocks: list[dict[str, Any]] | None = None,
) -> Any | None:
    """Fallback crop when no visual elements are detected near a caption."""
    text_blocks = text_blocks or []
    max_height = page_rect.height * 0.45
    min_height = 48

    above = page_rect.__class__(
        page_rect.x0 + margin,
        max(page_rect.y0, caption_rect.y0 - max_height),
        page_rect.x1 - margin,
        caption_rect.y0 - 4,
    )
    below = page_rect.__class__(
        page_rect.x0 + margin,
        caption_rect.y1 + 4,
        page_rect.x1 - margin,
        min(page_rect.y1, caption_rect.y1 + max_height),
    )

    def text_overlap_ratio(rect: Any) -> float:
        if not text_blocks:
            return 0.0
        overlap = sum(
            _rect_overlap_area(rect, block["bbox"]) for block in text_blocks
        )
        area = rect.width * rect.height
        return overlap / area if area > 0 else 1.0

    candidates = []
    for rect in (above, below):
        height = max(0, rect.y1 - rect.y0)
        if height < min_height:
            continue
        candidates.append((height, text_overlap_ratio(rect), rect))

    if not candidates:
        return None

    # Prefer the larger region, but drop a fallback that is mostly body text.
    candidates.sort(key=lambda x: (-x[0], x[1]))
    best = candidates[0]
    if best[1] > 0.5:
        return None
    return _clip_rect(best[2], page_rect)


def _cluster_rects(rects: list[Any], gap_threshold: float = 30) -> list[Any]:
    """Cluster rectangles that are close to each other."""
    if not rects:
        return []
    sorted_rects = sorted(rects, key=lambda r: (r.y0, r.x0))
    clusters: list[Any] = [sorted_rects[0]]
    for rect in sorted_rects[1:]:
        merged = False
        for idx, cluster in enumerate(clusters):
            if _rects_near(rect, cluster, gap_threshold):
                clusters[idx] = cluster | rect
                merged = True
                break
        if not merged:
            clusters.append(rect)
    return clusters


def _rects_near(a: Any, b: Any, gap_threshold: float) -> bool:
    """Return True if two rectangles are spatially close."""
    if a.x1 < b.x0:
        h_gap = b.x0 - a.x1
    elif b.x1 < a.x0:
        h_gap = a.x0 - b.x1
    else:
        h_gap = -min(a.x1, b.x1) + max(a.x0, b.x0)

    if a.y1 < b.y0:
        v_gap = b.y0 - a.y1
    elif b.y1 < a.y0:
        v_gap = a.y0 - b.y1
    else:
        v_gap = -min(a.y1, b.y1) + max(a.y0, b.y0)

    return h_gap <= gap_threshold and v_gap <= gap_threshold


def _split_subfigure_clusters(clusters: list[Any]) -> list[Any]:
    """Split a cluster into sub-clusters by internal whitespace gaps."""
    if len(clusters) <= 1:
        return clusters
    merged = clusters[0]
    for rect in clusters[1:]:
        merged |= rect
    if merged.width > merged.height * 1.4:
        rects = sorted(clusters, key=lambda r: r.x0)
        gaps = [
            (rects[i + 1].x0 - rects[i].x1, i)
            for i in range(len(rects) - 1)
        ]
        gaps.sort(reverse=True)
        best_gap, best_idx = gaps[0]
        if best_gap > merged.width * 0.08:
            left = _union_rects(rects[: best_idx + 1])
            right = _union_rects(rects[best_idx + 1 :])
            return [left, right]
    else:
        rects = sorted(clusters, key=lambda r: r.y0)
        gaps = [
            (rects[i + 1].y0 - rects[i].y1, i)
            for i in range(len(rects) - 1)
        ]
        gaps.sort(reverse=True)
        best_gap, best_idx = gaps[0]
        if best_gap > merged.height * 0.08:
            top = _union_rects(rects[: best_idx + 1])
            bottom = _union_rects(rects[best_idx + 1 :])
            return [top, bottom]
    return clusters


def _union_rects(rects: list[Any]) -> Any:
    if not rects:
        raise ValueError("cannot union empty rect list")
    result = rects[0]
    for rect in rects[1:]:
        result |= rect
    return result


def _extract_table_markdown(page: Any, rect: Any) -> str | None:
    """Try to extract a Markdown table from a PDF page region.

    First uses PyMuPDF's default line-based table finder; if no good table is
    found it falls back to the text-based strategy clipped to ``rect``.
    """
    best_table = None
    best_overlap = 0.0

    for strategy in ("lines", "text"):
        try:
            if strategy == "text":
                finder = page.find_tables(
                    clip=rect,
                    vertical_strategy="text",
                    horizontal_strategy="text",
                )
            else:
                finder = page.find_tables()
        except Exception:
            continue
        for table in finder:
            try:
                tbbox = table.bbox
            except Exception:
                continue
            overlap = _rect_overlap_area(rect, page.rect.__class__(*tbbox))
            if overlap > best_overlap:
                best_overlap = overlap
                best_table = table
        if best_table is not None and best_overlap > 0:
            break

    if best_table is None or best_overlap <= 0:
        return None

    try:
        rows = best_table.extract()
    except Exception:
        return None
    if not rows or len(rows) < 2:
        return None

    lines = []
    for i, row in enumerate(rows):
        cells = [str(cell or "").replace("|", "\\|").strip() for cell in row]
        lines.append("| " + " | ".join(cells) + " |")
        if i == 0:
            lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
    return "\n".join(lines)


def _rect_overlap_area(a: Any, b: Any) -> float:
    x_overlap = max(0, min(a.x1, b.x1) - max(a.x0, b.x0))
    y_overlap = max(0, min(a.y1, b.y1) - max(a.y0, b.y0))
    return x_overlap * y_overlap


def _clip_rect(rect: Any, page_rect: Any) -> Any:
    return rect & page_rect


def _split_paragraphs(text: str) -> list[str]:
    paragraphs = []
    current: list[str] = []
    for raw in text.splitlines():
        line = raw.strip()
        if not line:
            if current:
                paragraphs.append(_clean_text(" ".join(current)))
                current = []
            continue
        current.append(line)
    if current:
        paragraphs.append(_clean_text(" ".join(current)))
    return [paragraph for paragraph in paragraphs if paragraph]


def _clean_text(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()


def _safe_stem(path: Path) -> str:
    stem = re.sub(r"[^a-zA-Z0-9_-]+", "-", path.stem).strip("-").lower()
    return stem or "material"


def _label_asset_name(source_id: str, label: str, kind: str, fallback_index: int) -> str:
    """Create a stable asset name from a caption label, e.g. Figure 1 -> fig-1.png."""
    prefix = "table" if kind == "table" else "fig"
    match = re.search(r"(\d+)([a-zA-Z]?)\s*$", label)
    if match:
        number = match.group(1)
        suffix = match.group(2).lower()
        base = f"{source_id}-{prefix}-{number}{suffix}"
    else:
        base = f"{source_id}-{prefix}-{fallback_index}"
    return f"{base}.png"


def _register_source(
    course_dir: Path, source: Path, material_dir: Path, stats: dict[str, Any]
) -> None:
    workspace = course_dir / ".learnloop"
    workspace.mkdir(exist_ok=True)
    inventory = workspace / "source_inventory.yaml"
    source_id = _safe_stem(source)
    existing = inventory.read_text(encoding="utf-8") if inventory.exists() else ""
    if re.search(rf"^\s*-\s+id:\s*['\"]?{re.escape(source_id)}['\"]?\s*$", existing, re.M):
        return

    if not existing.strip():
        existing = "sources:\n"
    elif not existing.endswith("\n"):
        existing += "\n"
    if "sources:" not in existing:
        existing = "sources:\n" + existing

    try:
        raw_path = source.relative_to(course_dir).as_posix()
    except ValueError:
        raw_path = str(source)
    try:
        material_path = (material_dir / "material.json").relative_to(course_dir).as_posix()
    except ValueError:
        material_path = str(material_dir / "material.json")

    entry = (
        f"  - id: {source_id}\n"
        "    type: local-material\n"
        f"    title: {json.dumps(source.name, ensure_ascii=False)}\n"
        f"    path: {json.dumps(raw_path, ensure_ascii=False)}\n"
        "    status: ingested\n"
        f"    material: {json.dumps(material_path, ensure_ascii=False)}\n"
        f"    chunks: {int(stats.get('chunks', 0))}\n"
        f"    figures: {int(stats.get('figures', 0))}\n"
    )
    inventory.write_text(existing + entry, encoding="utf-8")


def _write_jsonl(path: Path, rows: list[dict[str, Any]]) -> None:
    path.write_text(
        "".join(json.dumps(row, ensure_ascii=False) + "\n" for row in rows),
        encoding="utf-8",
    )


def _figures_markdown(figures: list[dict[str, Any]]) -> str:
    parts = ["# Extracted Figures\n\n"]
    for figure in figures:
        caption = figure.get("caption", "")
        label = figure.get("label", "Figure")
        asset = figure["asset"]
        alt = f"{label} from {figure['source']}"
        kind = figure.get("kind", "figure")
        source_note = f"本地 {figure['source']} p.{figure['page']}"
        parts.append(
            "::: figure\n"
            f"src: {asset}\n"
            f"alt: {alt}\n"
            f"caption: {caption}\n"
            f"source: {source_note}\n"
            ":::\n\n"
        )
        if kind == "table":
            table_md = figure.get("table_markdown", "")
            if table_md:
                parts.append("**Extracted table**:\n\n")
                parts.append(table_md + "\n\n")
    return "".join(parts)
